

import tempfile

from io import BytesIO
from pathlib import Path
from typing import Any, ClassVar, cast

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..medias import Image
from ._iwork import extract_iwork_low_fidelity
from ._legacy_office import convert_legacy_office_bytes, infer_legacy_ole_kind
from .base import LogicalPageDocumentMixin, LLMDocumentPart
from ._utils import _table_to_markdown
from .pdf import PDF


def _convert_pptx_to_pdf_via_aspose(data: bytes) -> bytes | None:
    try:
        import aspose.slides as slides  # type: ignore
    except Exception:
        return None

    with tempfile.TemporaryDirectory(prefix='proj_pptx_pdf_') as tmp_dir:
        tmp_path = Path(tmp_dir)
        input_path = tmp_path / 'input.pptx'
        output_path = tmp_path / 'output.pdf'
        input_path.write_bytes(data)
        try:
            with slides.Presentation(str(input_path)) as presentation:
                presentation.save(str(output_path), slides.export.SaveFormat.PDF)
        except Exception:
            return None
        if output_path.exists() and output_path.stat().st_size > 0:
            return output_path.read_bytes()
    return None


class PPT(LogicalPageDocumentMixin):
    '''PPT 文档模型。

    NOTE: `.pptx` 走结构化解析；`.ppt` / `.key` 会在传入 LLM 时走尽力而为的低保真
    提取或临时转换，但 `model_dump` / `model_validate` 仍保持原始 source 数据不变。
    '''

    Abstract: ClassVar[bool] = False
    Type: ClassVar[str] = 'ppt'
    TypeNames: ClassVar[tuple[str, ...]] = ('presentation', 'ppt', 'pptx', 'key')
    Suffixes: ClassVar[tuple[str, ...]] = ('.ppt', '.pptx', '.key')
    MimeTypes: ClassVar[tuple[str, ...]] = ('application/vnd.openxmlformats-officedocument.presentationml.presentation',)

    def __init__(self, source: str | Path | bytes):
        super().__init__(source)
        self._page_cache: list[list[LLMDocumentPart]] | None = None

    def _validate_source(self) -> None:
        ...

    def _is_supported_format(self) -> bool:
        return self.suffix in {'.pptx'} or any(name.lower().startswith('ppt/') for name in self._zip_file_names())

    def _prepend_note(self, pages: list[list[LLMDocumentPart]], note: str) -> list[list[LLMDocumentPart]]:
        note_text = str(note or '').strip()
        if not note_text:
            return pages
        if not pages:
            return [[note_text]]
        return [[note_text, *pages[0]], *pages[1:]]

    def _assemble_low_fidelity_pages(
        self,
        *,
        note: str,
        text_blocks: list[str] | None = None,
        page_text_blocks: list[list[str]] | None = None,
        image_blobs: list[bytes] | None = None,
        preview_pdf: bytes | None = None,
    ) -> list[list[LLMDocumentPart]]:
        pages: list[list[LLMDocumentPart]]
        if preview_pdf:
            pages = cast(list[list[LLMDocumentPart]], PDF.Load(preview_pdf).extract_page_contents())
        else:
            pages = [[]]

        normalized_page_blocks = [
            [text.strip() for text in blocks if str(text).strip()]
            for blocks in (page_text_blocks or [])
        ]
        normalized_page_blocks = [blocks for blocks in normalized_page_blocks if blocks]

        if normalized_page_blocks:
            if not pages:
                pages = [list(blocks) for blocks in normalized_page_blocks]
            else:
                for page_index, blocks in enumerate(normalized_page_blocks):
                    block_items: list[LLMDocumentPart] = list(blocks)
                    if page_index < len(pages):
                        pages[page_index] = [*block_items, *pages[page_index]]
                    else:
                        pages.append(block_items)

        if text_blocks:
            text = '\n\n'.join(block.strip() for block in text_blocks if str(block).strip()).strip()
            if text:
                if not pages:
                    pages = [[text]]
                elif not normalized_page_blocks and pages[0]:
                    pages[0] = [text, *pages[0]]
                elif not normalized_page_blocks:
                    pages[0].append(text)

        for blob in image_blobs or []:
            try:
                image = Image(blob)
            except Exception:
                continue
            if not pages:
                pages = [[image]]
            else:
                pages[0].append(image)

        normalized_pages = [page for page in pages if page]
        return self._prepend_note(normalized_pages or [['(Empty presentation)']], note)

    def _extract_legacy_ppt_contents(self, source_suffix: str) -> list[list[LLMDocumentPart]]:
        result = convert_legacy_office_bytes(self.to_bytes(), kind='ppt', source_suffix=source_suffix)
        note = f'NOTE: The original {source_suffix or ".ppt"} file was extracted in low-fidelity mode for LLM input via {result.backend}.'
        if result.output_kind == 'pptx' and result.converted_bytes:
            pages = PPT.Load(result.converted_bytes).extract_page_contents()
            return self._prepend_note(pages, note)
        if result.output_kind == 'pdf' and result.converted_bytes:
            pages = PDF.Load(result.converted_bytes).extract_page_contents()
            return self._prepend_note(pages, note)  # type: ignore
        if result.text_content:
            return [[note, result.text_content]]
        raise ValueError('; '.join(result.warnings) or f'Failed to extract legacy presentation source: {source_suffix}.')

    def _extract_key_contents(self) -> list[list[LLMDocumentPart]]:
        result = extract_iwork_low_fidelity(self.to_bytes(), 'key')
        note = 'NOTE: The original .key file was extracted in low-fidelity mode for LLM input.'
        return self._assemble_low_fidelity_pages(
            note=note,
            text_blocks=result.text_blocks,
            page_text_blocks=result.page_text_blocks,
            image_blobs=result.image_blobs,
            preview_pdf=result.preview_pdf,
        )

    def to_pdf(self) -> bytes:
        if self.suffix == '.key':
            result = extract_iwork_low_fidelity(self.to_bytes(), 'key')
            if result.preview_pdf:
                return result.preview_pdf
            raise ValueError('No PDF preview backend available for .key source.')

        if self.suffix == '.ppt' or infer_legacy_ole_kind(self.to_bytes()) == 'ppt':
            result = convert_legacy_office_bytes(self.to_bytes(), kind='ppt', source_suffix=self.suffix or '.ppt')
            if result.output_kind == 'pdf' and result.converted_bytes:
                return result.converted_bytes
            if result.output_kind == 'pptx' and result.converted_bytes:
                return PPT.Load(result.converted_bytes).to_pdf()
            raise ValueError('; '.join(result.warnings) or 'No PDF conversion backend available for legacy PPT source.')

        if not self._is_supported_format():
            raise ValueError('Unsupported PPT source: supported formats are .ppt/.pptx/.key.')

        pdf_bytes = _convert_pptx_to_pdf_via_aspose(self.to_bytes())
        if pdf_bytes:
            return pdf_bytes
        raise ValueError('No PPTX-to-PDF backend available. Install `aspose.slides` to enable PPTX preview conversion.')

    def extract_page_contents(self, **kwargs: Any) -> list[list[LLMDocumentPart]]:
        if self._page_cache is not None:
            return self._page_cache
        if self.suffix == '.key':
            self._page_cache = self._extract_key_contents()
            return self._page_cache
        if self.suffix == '.ppt' or infer_legacy_ole_kind(self.to_bytes()) == 'ppt':
            self._page_cache = self._extract_legacy_ppt_contents(self.suffix or '.ppt')
            return self._page_cache
        if not self._is_supported_format():
            raise ValueError('Unsupported PPT source: supported formats are .ppt/.pptx/.key.')

        def _extract_shape_items(shape: Any) -> list[LLMDocumentPart]:
            items: list[LLMDocumentPart] = []
            shape_type = getattr(shape, 'shape_type', None)

            if shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in getattr(shape, 'shapes', []):
                    items.extend(_extract_shape_items(sub_shape))
                return items

            if getattr(shape, 'has_text_frame', False):
                text = str(getattr(shape, 'text', '') or '').strip()
                if text:
                    items.append(text)

            if getattr(shape, 'has_table', False):
                rows: list[list[str]] = []
                table = shape.table
                for row in table.rows:
                    rows.append([str(cell.text or '').strip() for cell in row.cells])
                table_md = _table_to_markdown(rows)
                if table_md:
                    items.append(table_md)

            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    if blob:
                        items.append(Image(blob))
                except Exception:
                    ...

            return items

        prs = Presentation(BytesIO(self.to_bytes()))
        pages: list[list[LLMDocumentPart]] = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            page_items: list[LLMDocumentPart] = []
            for shape in slide.shapes:
                page_items.extend(_extract_shape_items(shape))
            try:
                notes_frame = getattr(slide.notes_slide, 'notes_text_frame', None)
                notes_text = str(getattr(notes_frame, 'text', '') or '').strip()
            except Exception:
                notes_text = ''
            if notes_text:
                page_items.append(notes_text)
            if not page_items:
                page_items.append(f'(Empty slide {slide_index})')
            pages.append(page_items)

        self._page_cache = pages or [['(Empty presentation)']]
        return self._page_cache


__all__ = ['PPT']
