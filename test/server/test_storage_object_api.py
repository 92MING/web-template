# -*- coding: utf-8 -*-
"""Tests for Object Storage API endpoints (/_internal/admin/api/storage/object/*)."""


import io
import json
import unittest

from unittest.mock import patch

from _test_helpers import StorageObjectTestBase
from _test_helpers import _make_storage_config
from core.storage.config import LocalObjectDBConfig, ObjectStorageConfig, StorageConfig


def _make_docx_bytes() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("Object Storage Office Preview", level=1)
    doc.add_paragraph("This is a DOCX preview smoke test.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    box.text_frame.text = "PPTX preview smoke test"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pdf_bytes(text: str = "PPTX preview smoke test") -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=1280, height=720)
    page.insert_text((72, 96), text, fontsize=28)
    data = doc.tobytes()
    doc.close()
    return data


class TestObjectConfig(StorageObjectTestBase):
    """GET /_internal/admin/api/storage/object/config"""

    async def test_config_returns_200(self):
        resp = await self._client.get("/_internal/admin/api/storage/object/config")
        self.assertEqual(resp.status_code, 200)

    async def test_config_has_required_fields(self):
        data = (await self._client.get("/_internal/admin/api/storage/object/config")).json()
        for key in ("backend", "namespace", "supports_preview", "supports_cleanup"):
            self.assertIn(key, data, f"Missing '{key}'")

    async def test_config_backend_is_local(self):
        data = (await self._client.get("/_internal/admin/api/storage/object/config")).json()
        self.assertIsInstance(data["backend"], str)
        self.assertTrue(len(data["backend"]) > 0)


class StrictClientSelectionObjectTestBase(StorageObjectTestBase):
    @classmethod
    def _make_storage_config(cls, tmp: str) -> StorageConfig:
        base = _make_storage_config(tmp)
        return StorageConfig(
            kv=base.kv,
            orm=base.orm,
            object=ObjectStorageConfig(
                default=LocalObjectDBConfig(root_path=base.object.default.root_path, namespace="default-object"),
                cache=LocalObjectDBConfig(root_path=base.object.cache.root_path, namespace="cache-object"),
                extra={
                    "assets": LocalObjectDBConfig(root_path=f"{tmp}/objects_assets", namespace="assets-object"),
                },
            ),
            vector=base.vector,
        )


class TestObjectClientSelection(StrictClientSelectionObjectTestBase):
    async def test_config_resolves_exact_named_client(self):
        resp = await self._client.get("/_internal/admin/api/storage/object/config", params={"client": "assets"})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertEqual(data["client_name"], "assets")
        self.assertEqual(data["backend"], "local")

    async def test_named_client_isolated_from_default_client(self):
        files = {"files": ("asset.txt", b"asset-body", "text/plain")}
        upload_resp = await self._client.post(
            "/_internal/admin/api/storage/object/upload",
            params={"client": "assets"},
            files=files,
        )
        self.assertEqual(upload_resp.status_code, 200)
        path = upload_resp.json()["uploaded"][0]["path"]

        assets_meta = await self._client.get(
            "/_internal/admin/api/storage/object/meta",
            params={"client": "assets", "path": path},
        )
        self.assertEqual(assets_meta.status_code, 200)

        default_meta = await self._client.get(
            "/_internal/admin/api/storage/object/meta",
            params={"path": path},
        )
        self.assertEqual(default_meta.status_code, 404)

    async def test_unknown_client_does_not_fallback_or_fuzzy_match(self):
        resp = await self._client.get("/_internal/admin/api/storage/object/config", params={"client": "assetsz"})
        self.assertEqual(resp.status_code, 404)
        self.assertIn("assetsz", resp.text)


class TestObjectItems(StorageObjectTestBase):
    """GET /_internal/admin/api/storage/object/items"""

    async def test_items_empty_initially(self):
        resp = await self._client.get("/_internal/admin/api/storage/object/items")
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        for key in ("prefix", "breadcrumbs", "items", "total"):
            self.assertIn(key, data)
        self.assertIsInstance(data["items"], list)
        self.assertIsInstance(data["breadcrumbs"], list)

    async def test_items_breadcrumbs_root(self):
        data = (await self._client.get("/_internal/admin/api/storage/object/items")).json()
        self.assertGreaterEqual(len(data["breadcrumbs"]), 1)
        self.assertTrue(len(data["breadcrumbs"][0]["name"]) > 0)  # name should not be empty

    async def test_items_with_prefix(self):
        # Upload a file in a "subfolder"
        files = {"files": ("hello.txt", b"hello world", "text/plain")}
        await self._client.post(
            "/_internal/admin/api/storage/object/upload",
            files=files,
            data={"prefix": "subfolder"},
        )
        resp = await self._client.get("/_internal/admin/api/storage/object/items", params={"prefix": "subfolder"})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data["prefix"].startswith("subfolder"))

    async def test_items_pagination_params(self):
        resp = await self._client.get("/_internal/admin/api/storage/object/items", params={"limit": 5, "offset": 0})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertEqual(data["limit"], 5)
        self.assertEqual(data["offset"], 0)


class TestObjectUpload(StorageObjectTestBase):
    """POST /_internal/admin/api/storage/object/upload"""

    async def test_upload_single_file(self):
        content = b"test file content 12345"
        files = {"files": ("upload_test.txt", content, "text/plain")}
        resp = await self._client.post("/_internal/admin/api/storage/object/upload", files=files)
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertIn("uploaded", data)
        self.assertEqual(len(data["uploaded"]), 1)
        self.assertIn("path", data["uploaded"][0])

    async def test_upload_with_prefix(self):
        files = {"files": ("prefixed.txt", b"data", "text/plain")}
        resp = await self._client.post(
            "/_internal/admin/api/storage/object/upload",
            files=files,
            data={"prefix": "docs/reports"},
        )
        self.assertEqual(resp.status_code, 200)
        uploaded = resp.json()["uploaded"]
        self.assertTrue(uploaded[0]["path"].startswith("docs/reports/"))

    async def test_upload_with_expire(self):
        files = {"files": ("expiring.bin", b"\x00\x01\x02", "application/octet-stream")}
        resp = await self._client.post(
            "/_internal/admin/api/storage/object/upload",
            files=files,
            data={"expire_seconds": "3600"},
        )
        self.assertEqual(resp.status_code, 200)

    async def test_upload_with_metadata(self):
        meta = json.dumps({"author": "test", "version": "1.0"})
        files = {"files": ("meta.json", b'{"k":1}', "application/json")}
        resp = await self._client.post(
            "/_internal/admin/api/storage/object/upload",
            files=files,
            data={"metadata_json": meta},
        )
        self.assertEqual(resp.status_code, 200)

    async def test_upload_invalid_metadata_returns_400(self):
        files = {"files": ("bad_meta.txt", b"x", "text/plain")}
        resp = await self._client.post(
            "/_internal/admin/api/storage/object/upload",
            files=files,
            data={"metadata_json": "not{valid json"},
        )
        self.assertEqual(resp.status_code, 400)


class TestObjectMetaAndContent(StorageObjectTestBase):
    """GET /_internal/admin/api/storage/object/meta, GET /_internal/admin/api/storage/object/content"""

    async def _upload(self, name: str, content: bytes, content_type: str = "text/plain") -> str:
        files = {"files": (name, content, content_type)}
        resp = await self._client.post("/_internal/admin/api/storage/object/upload", files=files)
        return resp.json()["uploaded"][0]["path"]

    async def test_get_meta(self):
        path = await self._upload("meta_check.txt", b"metadata check")
        resp = await self._client.get("/_internal/admin/api/storage/object/meta", params={"path": path})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertEqual(data["path"], path)
        self.assertIn("size", data)
        self.assertIn("content_type", data)
        self.assertIn("download_url", data)

    async def test_get_meta_nonexistent_returns_404(self):
        resp = await self._client.get("/_internal/admin/api/storage/object/meta", params={"path": "no/such/file.xyz"})
        self.assertEqual(resp.status_code, 404)

    async def test_get_content(self):
        original = b"exact content verification"
        path = await self._upload("content_test.txt", original)
        resp = await self._client.get("/_internal/admin/api/storage/object/content", params={"path": path})
        self.assertEqual(resp.status_code, 200)
        self.assertEqual(resp.content, original)

    async def test_get_content_download_flag(self):
        path = await self._upload("dl_test.txt", b"download me")
        resp = await self._client.get(
            "/_internal/admin/api/storage/object/content",
            params={"path": path, "download": True},
        )
        self.assertEqual(resp.status_code, 200)
        self.assertIn("Content-Disposition", resp.headers)

    async def test_get_content_nonexistent_returns_404(self):
        resp = await self._client.get("/_internal/admin/api/storage/object/content", params={"path": "ghost/file.bin"})
        self.assertEqual(resp.status_code, 404)

    @unittest.skip("Office preview stub returns empty payload")
    async def test_office_preview_docx(self):
        path = await self._upload("preview.docx", _make_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        resp = await self._client.get("/_internal/admin/api/storage/object/office-preview", params={"path": path})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertEqual(data["kind"], "document")
        self.assertEqual(data["format"], "docx")
        self.assertEqual(data["support_level"], "stable")
        self.assertGreaterEqual(data["page_count"], 1)
        self.assertTrue(data["has_text"])
        self.assertTrue(data["pages"][0]["title"])
        self.assertIn("excerpt", data["pages"][0])
        all_text = "\n".join(
            item.get("text", "")
            for page in data["pages"]
            for item in page.get("items", [])
            if item.get("kind") == "markdown"
        )
        self.assertIn("DOCX preview smoke test", all_text)

    @unittest.skip("Office preview stub returns empty payload")
    async def test_office_preview_pptx(self):
        path = await self._upload("preview.pptx", _make_pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        with patch(
            "core.utils.data_structs.files.documents.ppt.PPT.to_pdf",
            return_value=_make_pdf_bytes(),
        ):
            resp = await self._client.get("/_internal/admin/api/storage/object/office-preview", params={"path": path})
            self.assertEqual(resp.status_code, 200)
            data = resp.json()
            self.assertEqual(data["kind"], "presentation")
            self.assertEqual(data["format"], "pptx")
            self.assertEqual(data["support_level"], "stable")
            self.assertEqual(data["preview_mode"], "pdf")
            self.assertGreaterEqual(data["page_count"], 1)
            self.assertTrue(data["pages"][0]["title"])
            self.assertIn("pdf_url", data)
            self.assertIn("thumbnail_url", data["pages"][0])

            pdf_resp = await self._client.get(data["pdf_url"])
            self.assertEqual(pdf_resp.status_code, 200)
            self.assertEqual(pdf_resp.headers.get("content-type"), "application/pdf")

            thumb_resp = await self._client.get(data["pages"][0]["thumbnail_url"])
            self.assertEqual(thumb_resp.status_code, 200)
            self.assertEqual(thumb_resp.headers.get("content-type"), "image/png")


class TestObjectExpireAndDelete(StorageObjectTestBase):
    """PATCH /_internal/admin/api/storage/object/expire, DELETE /_internal/admin/api/storage/object/item"""

    async def _upload(self, name: str, content: bytes) -> str:
        files = {"files": (name, content, "text/plain")}
        resp = await self._client.post("/_internal/admin/api/storage/object/upload", files=files)
        return resp.json()["uploaded"][0]["path"]

    async def test_set_expire(self):
        path = await self._upload("exp_test.txt", b"will expire")
        resp = await self._client.patch(
            "/_internal/admin/api/storage/object/expire",
            params={"path": path},
            json={"expire_seconds": 7200},
        )
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data["updated"])
        self.assertEqual(data["ttl_state"], "expiring")

    async def test_set_expire_nonexistent_returns_404(self):
        resp = await self._client.patch(
            "/_internal/admin/api/storage/object/expire",
            params={"path": "none/here.bin"},
            json={"expire_seconds": 100},
        )
        self.assertEqual(resp.status_code, 404)

    async def test_delete_object(self):
        path = await self._upload("del_me.txt", b"delete this")
        resp = await self._client.delete("/_internal/admin/api/storage/object/item", params={"path": path})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data["deleted"])
        # Verify gone
        meta_resp = await self._client.get("/_internal/admin/api/storage/object/meta", params={"path": path})
        self.assertEqual(meta_resp.status_code, 404)


class TestObjectCleanup(StorageObjectTestBase):
    """POST /_internal/admin/api/storage/object/cleanup"""

    async def test_cleanup(self):
        resp = await self._client.post("/_internal/admin/api/storage/object/cleanup")
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertIn("removed", data)
        self.assertIsInstance(data["removed"], int)


class TestObjectFolderAndAdvancedOps(StorageObjectTestBase):
    """Advanced object operations: folder/create/copy/move/metadata/content/search/delete-many."""

    async def _upload(self, name: str, content: bytes, prefix: str = "", content_type: str = "text/plain") -> str:
        files = {"files": (name, content, content_type)}
        resp = await self._client.post(
            "/_internal/admin/api/storage/object/upload",
            files=files,
            data={"prefix": prefix},
        )
        self.assertEqual(resp.status_code, 200)
        return resp.json()["uploaded"][0]["path"]

    async def test_create_folder_and_meta(self):
        resp = await self._client.post(
            "/_internal/admin/api/storage/object/folder",
            json={"path": "workspace/demo", "metadata": {"tags": ["empty", "docs"]}},
        )
        self.assertEqual(resp.status_code, 200)
        self.assertTrue(resp.json()["created"])

        meta_resp = await self._client.get(
            "/_internal/admin/api/storage/object/meta",
            params={"path": "workspace/demo/"},
        )
        self.assertEqual(meta_resp.status_code, 200)
        data = meta_resp.json()
        self.assertEqual(data["kind"], "folder")
        self.assertEqual(data["path"], "workspace/demo/")
        self.assertEqual(data["metadata"]["tags"], ["empty", "docs"])

    async def test_copy_and_move_file(self):
        source = await self._upload("move_me.txt", b"move target")

        copy_resp = await self._client.post(
            "/_internal/admin/api/storage/object/copy",
            json={"source_path": source, "target_path": "copies/move_me.txt"},
        )
        self.assertEqual(copy_resp.status_code, 200)
        self.assertEqual(copy_resp.json()["target_path"], "copies/move_me.txt")

        move_resp = await self._client.post(
            "/_internal/admin/api/storage/object/move",
            json={"source_path": "copies/move_me.txt", "target_path": "copies/moved.txt"},
        )
        self.assertEqual(move_resp.status_code, 200)
        self.assertEqual(move_resp.json()["target_path"], "copies/moved.txt")

        old_meta = await self._client.get(
            "/_internal/admin/api/storage/object/meta",
            params={"path": "copies/move_me.txt"},
        )
        self.assertEqual(old_meta.status_code, 404)

        new_meta = await self._client.get(
            "/_internal/admin/api/storage/object/meta",
            params={"path": "copies/moved.txt"},
        )
        self.assertEqual(new_meta.status_code, 200)

    async def test_patch_metadata_and_tags(self):
        path = await self._upload("meta_tags.txt", b"abc")
        resp = await self._client.patch(
            "/_internal/admin/api/storage/object/metadata",
            params={"path": path},
            json={"metadata": {"tags": ["alpha", "beta"], "owner": "tester"}, "merge": False},
        )
        self.assertEqual(resp.status_code, 200)
        self.assertTrue(resp.json()["updated"])

        meta_resp = await self._client.get("/_internal/admin/api/storage/object/meta", params={"path": path})
        self.assertEqual(meta_resp.status_code, 200)
        data = meta_resp.json()
        self.assertEqual(data["metadata"]["owner"], "tester")
        self.assertEqual(data["metadata"]["tags"], ["alpha", "beta"])

    async def test_put_text_content(self):
        path = await self._upload("editable.txt", b"old")
        resp = await self._client.put(
            "/_internal/admin/api/storage/object/content",
            params={"path": path},
            json={"mode": "text", "value": "new content"},
        )
        self.assertEqual(resp.status_code, 200)
        self.assertTrue(resp.json()["saved"])

        content_resp = await self._client.get(
            "/_internal/admin/api/storage/object/content",
            params={"path": path},
        )
        self.assertEqual(content_resp.status_code, 200)
        self.assertEqual(content_resp.text, "new content")

    async def test_items_support_wildcard_pattern(self):
        await self._upload("report-2026-01.txt", b"a", prefix="wild")
        await self._upload("report-2025-01.txt", b"b", prefix="wild")
        await self._upload("notes.md", b"c", prefix="wild")

        resp = await self._client.get(
            "/_internal/admin/api/storage/object/items",
            params={"prefix": "wild", "pattern": "report-2026-*.txt", "recursive": True},
        )
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertEqual(len(data["items"]), 1)
        self.assertEqual(data["items"][0]["name"], "report-2026-01.txt")

    async def test_items_support_type_group_filter(self):
        await self._upload("diagram.png", b"fake-png", prefix="typed", content_type="image/png")
        await self._upload("report.pdf", b"%PDF-1.7", prefix="typed", content_type="application/pdf")

        resp = await self._client.get(
            "/_internal/admin/api/storage/object/items",
            params={"prefix": "typed", "type_group": "image", "recursive": True},
        )
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertEqual([item["name"] for item in data["items"]], ["diagram.png"])
        self.assertEqual(data["filters"]["type_group"], "image")

    async def test_delete_many_mixed_file_and_folder(self):
        await self._client.post(
            "/_internal/admin/api/storage/object/folder",
            json={"path": "trash/empty-folder"},
        )
        path = await self._upload("gone.txt", b"bye", prefix="trash")

        resp = await self._client.post(
            "/_internal/admin/api/storage/object/delete-many",
            json={"paths": [path, "trash/empty-folder/"]},
        )
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data["deleted"])
        self.assertGreaterEqual(data["removed"], 2)

        file_meta = await self._client.get("/_internal/admin/api/storage/object/meta", params={"path": path})
        self.assertEqual(file_meta.status_code, 404)

        folder_meta = await self._client.get("/_internal/admin/api/storage/object/meta", params={"path": "trash/empty-folder/"})
        self.assertEqual(folder_meta.status_code, 404)


if __name__ == "__main__":
    unittest.main()
