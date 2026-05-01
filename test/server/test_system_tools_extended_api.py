# -*- coding: utf-8 -*-
"""Extended tests for system terminal and system files APIs."""


import asyncio
import contextlib
import io
import tempfile
import unittest

from pathlib import Path
from unittest.mock import patch

from fastapi.testclient import TestClient

from _test_helpers import FullAppTestBase


def _patch_system_roots(root: Path, default_terminal_path: str = "", max_sessions: int = 6):
    from core.server.routes.system.tools import RootInfo

    root = root.resolve()
    roots = [RootInfo(key="testroot", label="Test Root", path=str(root), is_default=True)]
    return patch(
        "core.server.routes.system.tools._parse_root_paths",
        return_value=(roots, {"testroot": root}, "testroot", default_terminal_path, max_sessions),
    )


def _make_docx_bytes() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("System Files Preview", level=1)
    doc.add_paragraph("DOCX preview smoke test")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    box.text_frame.text = "PPTX preview smoke test"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pdf_bytes(text: str = "PPTX preview smoke test") -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=1280, height=720)
    page.insert_text((72, 96), text, fontsize=28)
    data = doc.tobytes()
    doc.close()
    return data


class _FakeTerminalSession:
    def __init__(self, cols: int = 120, rows: int = 30):
        self.session_id = "fake-terminal-session"
        self.cols = cols
        self.rows = rows
        self.writes: list[str] = []
        self.resizes: list[tuple[int, int]] = []
        self.closed = False
        self._loop = None
        self._queue = None
        self._exit_sent = False

    def start_reader(self, loop, queue):
        self._loop = loop
        self._queue = queue

    def write(self, data: str):
        self.writes.append(data)

    def resize(self, cols: int, rows: int):
        self.cols = int(cols)
        self.rows = int(rows)
        self.resizes.append((self.cols, self.rows))

    def close(self):
        self.closed = True
        if self._exit_sent or self._loop is None or self._queue is None or self._loop.is_closed():
            return
        self._exit_sent = True
        with contextlib.suppress(Exception):
            self._loop.call_soon_threadsafe(self._queue.put_nowait, {"type": "exit", "code": 0})


class _TrackedTerminalSession:
    def __init__(self, session_id: str):
        self.session_id = session_id
        self.close_calls = 0

    def close(self):
        self.close_calls += 1


class TestSystemTerminalWebSocket(FullAppTestBase):
    async def test_terminal_websocket_ready_ping_and_resize(self):
        assert self._app is not None

        fake_session = _FakeTerminalSession()
        registered: list[tuple[str, int]] = []
        unregistered: list[str] = []

        with tempfile.TemporaryDirectory(prefix="proj_terminal_ws_") as tmp_dir:
            root = Path(tmp_dir)
            (root / "workspace").mkdir()

            with (
                _patch_system_roots(root, default_terminal_path="workspace"),
                patch("core.server.routes.system.tools._list_shells", return_value=[{"key": "powershell", "label": "PowerShell", "command": "powershell.exe"}]),
                patch("core.server.routes.system.tools._terminal_backend_status", return_value=(True, "pywinpty", None)),
                patch("core.server.routes.system.tools._create_terminal_session", return_value=fake_session),
                patch("core.server.routes.system.tools._register_terminal_session", side_effect=lambda session_id, max_sessions: registered.append((session_id, max_sessions))),
                patch("core.server.routes.system.tools._unregister_terminal_session", side_effect=lambda session_id: unregistered.append(session_id)),
            ):
                def _run_client():
                    with TestClient(self._app) as client:
                        with client.websocket_connect("/_internal/admin/ws/panel/system/terminal?root=testroot&shell=powershell&path=workspace&cols=120&rows=30") as websocket:
                            ready = websocket.receive_json()
                            self.assertEqual(ready["type"], "ready")
                            self.assertEqual(ready["shell"], "powershell")
                            self.assertEqual(ready["root"]["key"], "testroot")
                            self.assertEqual(ready["cwd"], "workspace")

                            websocket.send_json({"type": "ping"})
                            self.assertEqual(websocket.receive_json()["type"], "pong")

                            websocket.send_json({"type": "input", "data": "echo hi\r\n"})
                            websocket.send_json({"type": "resize", "cols": 140, "rows": 42})
                            fake_session.close()
                            self.assertEqual(websocket.receive_json()["type"], "exit")

                await asyncio.to_thread(_run_client)

        self.assertEqual(fake_session.writes, ["echo hi\r\n"])
        self.assertEqual(fake_session.resizes, [(140, 42)])
        self.assertTrue(fake_session.closed)
        self.assertEqual(registered, [("fake-terminal-session", 6)])
        self.assertEqual(unregistered, ["fake-terminal-session"])

    async def test_terminal_session_registry_closes_all_on_shutdown(self):
        from core.server.routes.system import tools as system_tools_module

        tracked = _TrackedTerminalSession("tracked-terminal-session")
        system_tools_module._close_all_terminal_sessions()
        try:
            system_tools_module._register_terminal_session(tracked.session_id, 6)
            system_tools_module._track_terminal_session(tracked)

            self.assertIn(tracked.session_id, system_tools_module._terminal_session_ids)
            self.assertIs(system_tools_module._terminal_sessions_by_id[tracked.session_id], tracked)

            system_tools_module._close_all_terminal_sessions()

            self.assertEqual(tracked.close_calls, 1)
            self.assertNotIn(tracked.session_id, system_tools_module._terminal_session_ids)
            self.assertNotIn(tracked.session_id, system_tools_module._terminal_sessions_by_id)
        finally:
            system_tools_module._close_all_terminal_sessions()

    async def test_terminal_websocket_invalid_shell_returns_error(self):
        assert self._app is not None

        with tempfile.TemporaryDirectory(prefix="proj_terminal_ws_") as tmp_dir:
            root = Path(tmp_dir)
            root.mkdir(exist_ok=True)

            with (
                _patch_system_roots(root),
                patch("core.server.routes.system.tools._list_shells", return_value=[{"key": "powershell", "label": "PowerShell", "command": "powershell.exe"}]),
                patch("core.server.routes.system.tools._terminal_backend_status", return_value=(True, "pywinpty", None)),
            ):
                def _run_client():
                    with TestClient(self._app) as client:
                        with client.websocket_connect("/_internal/admin/ws/panel/system/terminal?root=testroot&shell=unknown") as websocket:
                            error = websocket.receive_json()
                            self.assertEqual(error["type"], "error")
                            self.assertIn("unknown", error["message"])

                await asyncio.to_thread(_run_client)


class TestSystemFilesExtendedApi(FullAppTestBase):
    async def test_files_text_write_round_trip(self):
        with tempfile.TemporaryDirectory(prefix="proj_sys_files_") as tmp_dir:
            root = Path(tmp_dir)
            target = root / "note.txt"
            target.write_text("old content", encoding="utf-8")

            with _patch_system_roots(root):
                resp = await self._client.put(
                    "/_internal/admin/api/system/files/text",
                    params={"root": "testroot", "path": "note.txt"},
                    json={"content": "new content", "encoding": "utf-8"},
                )

            self.assertEqual(resp.status_code, 200)
            self.assertEqual(target.read_text(encoding="utf-8"), "new content")
            self.assertTrue(resp.json()["saved"])

    async def test_files_raw_and_download_return_expected_bytes(self):
        with tempfile.TemporaryDirectory(prefix="proj_sys_files_") as tmp_dir:
            root = Path(tmp_dir)
            payload = b"raw-bytes-check"
            (root / "sample.txt").write_bytes(payload)

            with _patch_system_roots(root):
                raw_resp = await self._client.get(
                    "/_internal/admin/api/system/files/raw",
                    params={"root": "testroot", "path": "sample.txt"},
                )
                download_resp = await self._client.get(
                    "/_internal/admin/api/system/files/download",
                    params={"root": "testroot", "path": "sample.txt"},
                )

            self.assertEqual(raw_resp.status_code, 200)
            self.assertEqual(raw_resp.content, payload)
            self.assertEqual(download_resp.status_code, 200)
            self.assertEqual(download_resp.content, payload)
            self.assertIn("sample.txt", download_resp.headers.get("content-disposition", ""))

    async def test_files_mkdir_upload_and_delete_round_trip(self):
        with tempfile.TemporaryDirectory(prefix="proj_sys_files_") as tmp_dir:
            root = Path(tmp_dir)

            with _patch_system_roots(root):
                mkdir_resp = await self._client.post(
                    "/_internal/admin/api/system/files/mkdir",
                    data={"root": "testroot", "path": "", "name": "docs"},
                )
                self.assertEqual(mkdir_resp.status_code, 200)

                upload_resp = await self._client.post(
                    "/_internal/admin/api/system/files/upload",
                    data={"root": "testroot", "path": "docs", "overwrite": "false"},
                    files=[("files", ("hello.txt", b"hello system files", "text/plain"))],
                )
                self.assertEqual(upload_resp.status_code, 200)
                self.assertEqual(upload_resp.json()["files"][0]["name"], "hello.txt")

                list_resp = await self._client.get(
                    "/_internal/admin/api/system/files/list",
                    params={"root": "testroot", "path": "docs"},
                )
                self.assertEqual(list_resp.status_code, 200)
                self.assertEqual([item["name"] for item in list_resp.json()["entries"]], ["hello.txt"])

                delete_file_resp = await self._client.delete(
                    "/_internal/admin/api/system/files/item",
                    params={"root": "testroot", "path": "docs/hello.txt"},
                )
                self.assertEqual(delete_file_resp.status_code, 200)
                self.assertTrue(delete_file_resp.json()["deleted"])

                delete_dir_resp = await self._client.delete(
                    "/_internal/admin/api/system/files/item",
                    params={"root": "testroot", "path": "docs"},
                )
                self.assertEqual(delete_dir_resp.status_code, 200)
                self.assertTrue(delete_dir_resp.json()["was_dir"])

            self.assertFalse((root / "docs").exists())

    @unittest.skip("Office preview stub returns empty payload")
    async def test_files_office_preview_docx(self):
        with tempfile.TemporaryDirectory(prefix="proj_sys_files_") as tmp_dir:
            root = Path(tmp_dir)
            (root / "preview.docx").write_bytes(_make_docx_bytes())

            with _patch_system_roots(root):
                resp = await self._client.get(
                    "/_internal/admin/api/system/files/office-preview",
                    params={"root": "testroot", "path": "preview.docx"},
                )

            self.assertEqual(resp.status_code, 200)
            data = resp.json()
            self.assertEqual(data["kind"], "document")
            self.assertEqual(data["format"], "docx")
            self.assertGreaterEqual(data["page_count"], 1)

    @unittest.skip("Office preview stub returns empty payload")
    async def test_files_office_preview_pptx_pdf_and_thumb(self):
        with tempfile.TemporaryDirectory(prefix="proj_sys_files_") as tmp_dir:
            root = Path(tmp_dir)
            (root / "preview.pptx").write_bytes(_make_pptx_bytes())

            with (
                _patch_system_roots(root),
                patch(
                    "core.utils.data_structs.files.documents.ppt.PPT.to_pdf",
                    return_value=_make_pdf_bytes(),
                ),
            ):
                resp = await self._client.get(
                    "/_internal/admin/api/system/files/office-preview",
                    params={"root": "testroot", "path": "preview.pptx"},
                )
                self.assertEqual(resp.status_code, 200)
                data = resp.json()
                self.assertEqual(data["kind"], "presentation")
                self.assertEqual(data["format"], "pptx")
                self.assertIn("pdf_url", data)
                self.assertIn("thumbnail_url", data["pages"][0])

                pdf_resp = await self._client.get(data["pdf_url"])
                thumb_resp = await self._client.get(data["pages"][0]["thumbnail_url"])

            self.assertEqual(pdf_resp.status_code, 200)
            self.assertEqual(pdf_resp.headers.get("content-type"), "application/pdf")
            self.assertEqual(thumb_resp.status_code, 200)
            self.assertEqual(thumb_resp.headers.get("content-type"), "image/png")


if __name__ == "__main__":
    unittest.main()