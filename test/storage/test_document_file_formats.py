

import asyncio
import base64
import json
import os
import sys
import tempfile
import zipfile

from io import BytesIO
from pathlib import Path
from typing import cast
from unittest import TestCase, main
from unittest.mock import patch

from pydantic import BaseModel
from PIL import Image as PILImage

_PROJECT_ROOT = Path(__file__).resolve().parents[2]
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

from core.utils.data_structs import CSV, Doc, Excel, File, HTML, Image, JSON, Markdown, PlainText, PPT, RTF, TOML, TSV, TXT, XML, YAML
from core.utils.data_structs.files.documents.base import LLMDocumentMixin
from core.utils.data_structs.files.documents._legacy_office import ConversionResult
from core.ai.completion import _normalize_msg_content_for_thinkthinksyn, _to_openai_content


class FilePayloadModel(BaseModel):
    data: File


def _await(value):
    return asyncio.run(value)


def _make_image_png_bytes(width: int = 64, height: int = 64) -> bytes:
    img = PILImage.new('RGB', (width, height), color=(80, 180, 120))
    buf = BytesIO()
    img.save(buf, format='PNG')
    return buf.getvalue()


def _make_docx_bytes(image_png_bytes: bytes) -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading('Doc Title', level=1)
    doc.add_paragraph('DOC_TEXT_1')
    doc.add_picture(BytesIO(image_png_bytes))
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = 'A'
    table.cell(0, 1).text = 'B'
    table.cell(1, 0).text = '1'
    table.cell(1, 1).text = '2'

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes(image_png_bytes: bytes) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    if title is not None:
        title.text = 'Slide Title'
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
    textbox.text_frame.text = 'PPT_TEXT_1'
    slide.shapes.add_picture(BytesIO(image_png_bytes), Inches(1), Inches(2.2), width=Inches(2))

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes() -> bytes:
    from openpyxl import Workbook
    from openpyxl.worksheet.worksheet import Worksheet

    wb = Workbook()
    ws = cast(Worksheet, wb.active)
    ws.title = 'Alpha'
    ws.append(['Name', 'Score'])
    ws.append(['Alice', 90])
    ws2 = wb.create_sheet('Beta')
    ws2.append(['Item', 'Qty'])
    ws2.append(['Pen', 3])

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_preview_pdf_bytes(text: str, image_png_bytes: bytes | None = None) -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=320, height=240)
    page.insert_text((24, 40), text)
    if image_png_bytes:
        page.insert_image(fitz.Rect(24, 60, 140, 140), stream=image_png_bytes)
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def _make_preview_pdf_pages(texts: list[str], image_png_bytes: bytes | None = None) -> bytes:
    import fitz

    doc = fitz.open()
    for text in texts:
        page = doc.new_page(width=320, height=240)
        page.insert_text((24, 40), text)
        if image_png_bytes:
            page.insert_image(fitz.Rect(24, 60, 140, 140), stream=image_png_bytes)
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def _make_pages_bytes(preview_pdf_bytes: bytes, image_png_bytes: bytes) -> bytes:
    buf = BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        zf.writestr('QuickLook/Preview.pdf', preview_pdf_bytes)
        zf.writestr('preview.jpg', image_png_bytes)
        zf.writestr('index.xml', '<document><page><p>PAGES_XML_TEXT_1</p></page><page><p>PAGES_XML_TEXT_2</p></page></document>')
    return buf.getvalue()


def _make_key_bytes(preview_pdf_bytes: bytes, image_png_bytes: bytes) -> bytes:
    buf = BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        zf.writestr('QuickLook/Preview.pdf', preview_pdf_bytes)
        zf.writestr('preview.jpg', image_png_bytes)
        zf.writestr('index.apxl', '<presentation><slide><text>KEY_XML_TEXT_1</text></slide><slide><text>KEY_XML_TEXT_2</text></slide></presentation>')
    return buf.getvalue()


def _make_hwpx_bytes(image_png_bytes: bytes) -> bytes:
    buf = BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        zf.writestr('mimetype', 'application/hwp+zip')
        zf.writestr('Preview/PrvText.txt', 'HWPX_PREVIEW_TEXT')
        zf.writestr('Contents/section0.xml', '<hp:section xmlns:hp="urn:hwpx"><hp:p><hp:t>HWPX_SECTION_TEXT_1</hp:t></hp:p></hp:section>')
        zf.writestr('Contents/section1.xml', '<hp:section xmlns:hp="urn:hwpx"><hp:p><hp:t>HWPX_SECTION_TEXT_2</hp:t></hp:p></hp:section>')
        zf.writestr('BinData/BIN0001.png', image_png_bytes)
    return buf.getvalue()


def _write_temp_file(suffix: str, data: bytes) -> Path:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    try:
        tmp.write(data)
        tmp.flush()
    finally:
        tmp.close()
    return Path(tmp.name)


def _make_csv_bytes() -> bytes:
    return b'name,score\nAlice,95\nBob,88\n'


def _make_tsv_bytes() -> bytes:
    return 'name\tscore\nAlice\t95\nBob\t88\n'.encode('utf-8')


def _make_odt_bytes() -> bytes:
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    doc = OpenDocumentText()
    doc.text.addElement(P(text='ODT_TEXT_1'))
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_ods_bytes() -> bytes:
    from odf.opendocument import OpenDocumentSpreadsheet
    from odf.table import Table, TableCell, TableRow
    from odf.text import P

    doc = OpenDocumentSpreadsheet()
    sheet = Table(name='Gamma')

    header = TableRow()
    for value in ('Name', 'Qty'):
        cell = TableCell()
        cell.addElement(P(text=value))
        header.addElement(cell)
    sheet.addElement(header)

    row = TableRow()
    for value in ('Book', '5'):
        cell = TableCell()
        cell.addElement(P(text=value))
        row.addElement(cell)
    sheet.addElement(row)

    doc.spreadsheet.addElement(sheet)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


class DocumentFileTest(TestCase):
    @classmethod
    def setUpClass(cls):
        cls.image_png_bytes = _make_image_png_bytes()
        cls.doc = Doc.Load(_make_docx_bytes(cls.image_png_bytes))
        cls.odt_doc = Doc.Load(_make_odt_bytes())
        cls.ppt = PPT.Load(_make_pptx_bytes(cls.image_png_bytes))
        cls.excel = Excel.Load(_make_xlsx_bytes())
        cls.ods_excel = Excel.Load(_make_ods_bytes())
        cls.html = HTML.Load(
            (
                '<html><body><h1>HTML</h1><p>'
                + ('hello ' * 300)
                + '</p><img src="data:image/png;base64,'
                + base64.b64encode(cls.image_png_bytes).decode('utf-8')
                + '"></body></html>'
            ).encode('utf-8')
        )
        cls.rtf = RTF.Load(r'{\rtf1\ansi RTF_TEXT}'.encode('utf-8'))
        cls.json = JSON.Load(b"{'name': 'demo', // comment\n 'items': [1, 2, 3]}\n")
        cls.yaml = YAML.Load(b'name: demo\nitems:\n  - 1\n  - 2\n')
        cls.toml = TOML.Load(b'title = "demo"\n[owner]\nname = "Bot"\n')
        cls.xml = XML.Load(b'<root><item>XML_1</item><item>XML_2</item></root>')
        cls.csv = CSV.Load(_make_csv_bytes())
        cls.tsv = TSV.Load(_make_tsv_bytes())
        cls.txt = TXT.Load('TXT_TEXT_1\nTXT_TEXT_2'.encode('utf-8'))
        cls.plaintext = PlainText.Load('PLAINTEXT_1\nPLAINTEXT_2'.encode('utf-8'))
        cls.markdown = Markdown.Load(
            (
                '---\n'
                'title: Demo Markdown\n'
                'tags:\n'
                '  - demo\n'
                '  - parser\n'
                '---\n\n'
                '# Demo Markdown\n\n'
                '- item 1\n'
                '- item 2\n\n'
                '## Details\n\n'
                'This is **important** markdown content.\n'
            ).encode('utf-8')
        )

    def test_file_roundtrip_for_new_documents(self):
        for item, expected_type in (
            (self.doc, 'doc'),
            (self.ppt, 'ppt'),
            (self.excel, 'excel'),
            (self.html, 'html'),
            (self.rtf, 'rtf'),
        ):
            model = FilePayloadModel(data=item)  # type: ignore[arg-type]
            dumped = json.loads(model.model_dump_json())
            self.assertEqual(dumped['data']['type'], expected_type)
            roundtrip = FilePayloadModel.model_validate_json(model.model_dump_json())
            self.assertIsInstance(roundtrip.data, type(item))

    def test_doc_to_llm_contains_text_and_image(self):
        parts = list(_await(self.doc.to_llm()))
        self.assertTrue(any(isinstance(part, str) and 'DOC_TEXT_1' in part for part in parts))
        self.assertTrue(any(isinstance(part, Image) for part in parts))

    def test_doc_image_mode_and_suffixless_odt_are_supported(self):
        image_parts = list(_await(self.doc.to_llm(mode='image')))
        self.assertTrue(image_parts)
        self.assertTrue(all(isinstance(part, Image) for part in image_parts))

        odt_parts = list(_await(self.odt_doc.to_llm()))
        self.assertTrue(any(isinstance(part, str) and 'ODT_TEXT_1' in part for part in odt_parts))

    def test_ppt_to_llm_contains_text_and_image(self):
        parts = list(_await(self.ppt.to_llm()))
        self.assertTrue(any(isinstance(part, str) and 'PPT_TEXT_1' in part for part in parts))
        self.assertTrue(any(isinstance(part, Image) for part in parts))

    def test_ppt_image_mode_returns_images(self):
        image_parts = list(_await(self.ppt.to_llm(mode='image')))
        self.assertTrue(image_parts)
        self.assertTrue(all(isinstance(part, Image) for part in image_parts))

    def test_html_to_llm_converts_long_html_and_extracts_media(self):
        parts = list(_await(self.html.to_llm()))
        self.assertTrue(any(isinstance(part, str) and 'converted from original HTML to Markdown' in part for part in parts))
        self.assertTrue(any(isinstance(part, Image) for part in parts))

    def test_rtf_short_content_returns_raw_source(self):
        parts = list(_await(self.rtf.to_llm()))
        self.assertEqual(1, len(parts))
        text = cast(str, parts[0])
        self.assertIsInstance(text, str)
        self.assertIn('RTF_TEXT', text)

    def test_excel_to_llm_marks_sheet_names(self):
        parts = list(_await(self.excel.to_llm()))
        self.assertEqual(1, len(parts))
        text = cast(str, parts[0])
        self.assertIn('#Alpha', text)
        self.assertIn('#Beta', text)
        self.assertIn('Alice\t90', text)

    def test_suffixless_ods_excel_is_supported(self):
        parts = list(_await(self.ods_excel.to_llm()))
        self.assertEqual(1, len(parts))
        text = cast(str, parts[0])
        self.assertIn('#Gamma', text)
        self.assertIn('Book\t5', text)

    def test_completion_normalizers_expand_new_document_types(self):
        openai_content = _await(_to_openai_content(self.html))
        self.assertIsInstance(openai_content, list)
        self.assertTrue(any(isinstance(part, dict) and part.get('type') == 'text' for part in cast(list[dict[str, object]], openai_content)))
        self.assertTrue(any(isinstance(part, dict) and part.get('type') == 'image_url' for part in cast(list[dict[str, object]], openai_content)))

        text, medias = _await(_normalize_msg_content_for_thinkthinksyn(self.ppt))
        self.assertIn('PPT_TEXT_1', text)
        self.assertIsInstance(medias, dict)
        self.assertTrue(bool(medias))

    def test_pages_low_fidelity_support_preserves_original_bytes(self):
        pages_bytes = _make_pages_bytes(_make_preview_pdf_pages(['PAGES_PREVIEW_TEXT_1', 'PAGES_PREVIEW_TEXT_2'], self.image_png_bytes), self.image_png_bytes)
        path = _write_temp_file('.pages', pages_bytes)
        try:
            doc = Doc.Load(path)
            parts = list(_await(doc.to_llm()))
            self.assertGreaterEqual(doc.page_count, 2)
            self.assertTrue(any(isinstance(part, str) and 'PAGES_XML_TEXT_2' in part for part in parts))
            self.assertTrue(any(isinstance(part, Image) for part in parts))

            model = FilePayloadModel(data=doc)
            roundtrip = FilePayloadModel.model_validate_json(model.model_dump_json())
            self.assertIsInstance(roundtrip.data, Doc)
            self.assertEqual(roundtrip.data.to_bytes(), pages_bytes)
        finally:
            path.unlink(missing_ok=True)

    def test_key_low_fidelity_support_preserves_original_bytes(self):
        key_bytes = _make_key_bytes(_make_preview_pdf_pages(['KEY_PREVIEW_TEXT_1', 'KEY_PREVIEW_TEXT_2'], self.image_png_bytes), self.image_png_bytes)
        path = _write_temp_file('.key', key_bytes)
        try:
            ppt = PPT.Load(path)
            parts = list(_await(ppt.to_llm()))
            self.assertGreaterEqual(ppt.page_count, 2)
            self.assertTrue(any(isinstance(part, str) and 'KEY_XML_TEXT_2' in part for part in parts))
            self.assertTrue(any(isinstance(part, Image) for part in parts))

            model = FilePayloadModel(data=ppt)
            roundtrip = FilePayloadModel.model_validate_json(model.model_dump_json())
            self.assertIsInstance(roundtrip.data, PPT)
            self.assertEqual(roundtrip.data.to_bytes(), key_bytes)
        finally:
            path.unlink(missing_ok=True)

    def test_hwpx_low_fidelity_support(self):
        hwpx = Doc.Load(_make_hwpx_bytes(self.image_png_bytes))
        parts = list(_await(hwpx.to_llm()))
        self.assertGreaterEqual(hwpx.page_count, 2)
        self.assertTrue(any(isinstance(part, str) and 'HWPX_SECTION_TEXT_2' in part for part in parts))
        self.assertTrue(any(isinstance(part, Image) for part in parts))

    def test_hwp_and_wps_fallback_text_support(self):
        for suffix, expected in (('.hwp', 'HWP_TEXT_1'), ('.wps', 'WPS_TEXT_1')):
            path = _write_temp_file(suffix, expected.encode('utf-16le'))
            try:
                doc = Doc.Load(path)
                parts = list(_await(doc.to_llm()))
                self.assertTrue(any(isinstance(part, str) and expected in part for part in parts))
            finally:
                path.unlink(missing_ok=True)

    def test_legacy_doc_conversion_preserves_original_bytes(self):
        original_bytes = b'LEGACY_DOC_BYTES'
        path = _write_temp_file('.doc', original_bytes)
        try:
            with patch(
                'core.utils.data_structs.files.documents.doc.convert_legacy_office_bytes',
                return_value=ConversionResult(output_kind='docx', backend='mock-docx', converted_bytes=_make_docx_bytes(self.image_png_bytes)),
            ):
                doc = Doc.Load(path)
                parts = list(_await(doc.to_llm()))
                self.assertTrue(any(isinstance(part, str) and 'DOC_TEXT_1' in part for part in parts))

                model = FilePayloadModel(data=doc)
                roundtrip = FilePayloadModel.model_validate_json(model.model_dump_json())
                self.assertEqual(roundtrip.data.to_bytes(), original_bytes)
        finally:
            path.unlink(missing_ok=True)

    def test_legacy_ppt_conversion_preserves_original_bytes(self):
        original_bytes = b'LEGACY_PPT_BYTES'
        path = _write_temp_file('.ppt', original_bytes)
        try:
            with patch(
                'core.utils.data_structs.files.documents.ppt.convert_legacy_office_bytes',
                return_value=ConversionResult(output_kind='pptx', backend='mock-pptx', converted_bytes=_make_pptx_bytes(self.image_png_bytes)),
            ):
                ppt = PPT.Load(path)
                parts = list(_await(ppt.to_llm()))
                self.assertTrue(any(isinstance(part, str) and 'PPT_TEXT_1' in part for part in parts))

                model = FilePayloadModel(data=ppt)
                roundtrip = FilePayloadModel.model_validate_json(model.model_dump_json())
                self.assertEqual(roundtrip.data.to_bytes(), original_bytes)
        finally:
            path.unlink(missing_ok=True)

    def test_new_structured_document_classes(self):
        self.assertEqual(self.json.to_dict()['name'], 'demo')
        self.assertEqual(self.json.get_path('items[1]'), 2)
        self.assertIn('items[2]', self.json.to_flat_dict())
        self.assertEqual(self.yaml.to_dict()['name'], 'demo')
        self.assertEqual(self.toml.to_dict()['owner']['name'], 'Bot')
        self.assertIn('XML_2', '\n'.join(self.xml.extract_markup_blocks()))
        self.assertEqual(self.xml.find_texts('item'), ['XML_1', 'XML_2'])
        self.assertEqual(self.xml.tag_counts().get('item'), 2)
        self.assertEqual(self.csv.sheet_names(), ['Sheet1'])
        self.assertEqual(self.csv.infer_header_row(), 0)
        self.assertEqual(self.csv.to_records()[1]['score'], '88')
        self.assertEqual(self.tsv.to_dicts()[0]['name'], 'Alice')

        json_parts = list(_await(self.json.to_llm(flatten=True)))
        self.assertTrue(any('items[2]' in part for part in json_parts if isinstance(part, str)))

    def test_plain_text_family_and_unknown_text_fallback(self):
        self.assertIn('TXT_TEXT_1', _await(self.txt.to_llm())[0])
        self.assertEqual(self.txt.paragraphs(), ['TXT_TEXT_1\nTXT_TEXT_2'])
        self.assertIn('PLAINTEXT_2', _await(self.plaintext.to_llm())[0])
        self.assertEqual(self.plaintext.snippet(12), 'PLAINTEXT_1…')
        markdown_parts = _await(self.markdown.to_llm())
        self.assertIn('Parsed Markdown front matter', markdown_parts[0])
        self.assertTrue(any(isinstance(part, str) and '## Details' in part for part in markdown_parts))
        self.assertEqual(self.markdown.front_matter()['title'], 'Demo Markdown')
        self.assertEqual(self.markdown.heading_sections()[1]['title'], 'Details')
        markdown_text = self.markdown.to_plain_text().replace('\n', ' ')
        self.assertIn('important', markdown_text)
        self.assertIn('markdown content', markdown_text)

        path = _write_temp_file('', 'UNKNOWN_TEXT_1\nUNKNOWN_TEXT_2'.encode('utf-8'))
        try:
            model = FilePayloadModel.model_validate({'data': str(path)})
            self.assertIsInstance(model.data, PlainText)
            self.assertIn('UNKNOWN_TEXT_1', _await(model.data.to_llm())[0])
        finally:
            path.unlink(missing_ok=True)

    def test_unknown_binary_file_still_errors(self):
        path = _write_temp_file('', b'\x00\x01\x02\x03\x04\x05')
        try:
            with self.assertRaises(ValueError):
                FilePayloadModel.model_validate({'data': str(path)})
        finally:
            path.unlink(missing_ok=True)

    def test_document_type_metadata(self):
        self.assertTrue(issubclass(Doc, LLMDocumentMixin))
        self.assertEqual('doc', Doc.Type)
        self.assertIn('docx', Doc.TypeNames)
        self.assertIn('pages', Doc.TypeNames)
        self.assertIn('hwpx', Doc.TypeNames)
        self.assertEqual('ppt', PPT.Type)
        self.assertIn('pptx', PPT.TypeNames)
        self.assertIn('key', PPT.TypeNames)
        self.assertEqual('excel', Excel.Type)
        self.assertIn('xlsx', Excel.TypeNames)
        self.assertEqual('json', JSON.Type)
        self.assertEqual('yaml', YAML.Type)
        self.assertEqual('toml', TOML.Type)
        self.assertEqual('xml', XML.Type)
        self.assertEqual('csv', CSV.Type)
        self.assertEqual('tsv', TSV.Type)
        self.assertEqual('txt', TXT.Type)
        self.assertEqual('plaintext', PlainText.Type)


if __name__ == '__main__':
    main()
