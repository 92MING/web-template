

import os
import sys
from pathlib import Path

_curr_dir = os.path.dirname(os.path.abspath(__file__))
_PROJECT_ROOT = Path(__file__).resolve().parents[2]
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

import json
import base64
import wave
import tempfile
import fitz
import numpy as np

from io import BytesIO
from pathlib import Path
from typing import cast
from pydantic import BaseModel
from unittest import TestCase, main
from PIL import Image as PILImage
from openpyxl import Workbook
from openpyxl.worksheet.worksheet import Worksheet
from moviepy.video.VideoClip import ImageClip
from moviepy.audio.io.AudioFileClip import AudioFileClip
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from core.utils.data_structs import Audio, CSV, Doc, Excel, File, HTML, Image, JSON, Markdown, PDF, PlainText, PPT, RTF, TOML, TSV, TXT, Video, XML, YAML
from core.utils.data_structs.files.base import FileTypeMetaProtocol

_TEST_RESOURCES_DIR = _PROJECT_ROOT / 'resources' / 'test'

class MediaPayloadModel(BaseModel):
    audio: Audio
    image: Image
    video: Video
    pdf: PDF

class DocumentPayloadModel(BaseModel):
    pdf: PDF
    ppt: PPT
    html: HTML
    doc: Doc
    rtf: RTF
    excel: Excel

class StructuredDocumentPayloadModel(BaseModel):
    json_doc: JSON
    yaml_doc: YAML
    toml_doc: TOML
    xml_doc: XML
    csv_doc: CSV
    tsv_doc: TSV
    txt_doc: TXT
    plaintext_doc: PlainText
    markdown_doc: Markdown

class FilePayloadModel(BaseModel):
    data: File

def _make_audio_wav_bytes(duration_sec: float = 0.5, sample_rate: int = 16000) -> bytes:
    t = np.linspace(0, duration_sec, int(sample_rate * duration_sec), endpoint=False)
    wave_data = (0.2 * np.sin(2 * np.pi * 440 * t)).astype(np.float32)
    pcm16 = np.clip(wave_data * 32767, -32768, 32767).astype(np.int16)

    buf = BytesIO()
    with wave.open(buf, 'wb') as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(sample_rate)
        wf.writeframes(pcm16.tobytes())
    return buf.getvalue()

def _make_image_png_bytes(width: int = 64, height: int = 64) -> bytes:
    arr = np.zeros((height, width, 3), dtype=np.uint8)
    arr[:, :, 1] = 180
    img = PILImage.fromarray(arr, mode='RGB')
    buf = BytesIO()
    img.save(buf, format='PNG')
    return buf.getvalue()

def _make_video_mp4_bytes(image_png_bytes: bytes, audio_wav_bytes: bytes, duration_sec: float = 0.5) -> bytes:
    with tempfile.TemporaryDirectory(prefix='proj_file_test_') as td:
        tmp_dir = Path(td)
        image_path = tmp_dir / 'frame.png'
        audio_path = tmp_dir / 'audio.wav'
        video_path = tmp_dir / 'video.mp4'

        image_path.write_bytes(image_png_bytes)
        audio_path.write_bytes(audio_wav_bytes)

        arr = np.array(PILImage.open(BytesIO(image_png_bytes)).convert('RGB'))
        clip = ImageClip(arr)
        if hasattr(clip, 'with_duration'):
            clip = clip.with_duration(duration_sec)
        else:
            clip = clip.set_duration(duration_sec)  # type: ignore[attr-defined]

        audio_clip = AudioFileClip(str(audio_path))
        if hasattr(audio_clip, 'subclipped'):
            audio_clip = audio_clip.subclipped(0, duration_sec)
        else:
            audio_clip = audio_clip.subclip(0, duration_sec)  # type: ignore[attr-defined]

        if hasattr(clip, 'with_audio'):
            clip = clip.with_audio(audio_clip)
        else:
            clip = clip.set_audio(audio_clip)  # type: ignore[attr-defined]

        clip.write_videofile(
            str(video_path),
            codec='libx264',
            audio_codec='aac',
            fps=24,
            logger=None,
        )

        try:
            clip.close()
        except Exception:
            pass
        try:
            audio_clip.close()
        except Exception:
            pass

        return video_path.read_bytes()

def _make_pdf_bytes() -> bytes:
    doc = fitz.open()
    page = doc.new_page(width=300, height=200)
    page.insert_text((36, 72), 'File Test PDF')
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes

def _make_pdf_mixed_bytes(image_png_bytes: bytes) -> bytes:
    doc = fitz.open()
    page = doc.new_page(width=360, height=260)
    page.insert_text((36, 40), 'MIX_TEXT_1')
    page.insert_image(fitz.Rect(36, 60, 180, 150), stream=image_png_bytes)
    page.insert_text((36, 190), 'MIX_TEXT_2')
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes

def _make_docx_bytes(image_png_bytes: bytes) -> bytes:
    doc = Document()
    doc.add_heading('Doc Title', level=1)
    doc.add_paragraph('DOC_TEXT_1')
    doc.add_picture(BytesIO(image_png_bytes))
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = 'A'
    table.cell(0, 1).text = 'B'
    table.cell(1, 0).text = '1'
    table.cell(1, 1).text = '2'

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()

def _make_pptx_bytes(image_png_bytes: bytes) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    if title is not None:
        title.text = 'Slide Title'
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
    textbox.text_frame.text = 'PPT_TEXT_1'
    slide.shapes.add_picture(BytesIO(image_png_bytes), Inches(1), Inches(2.2), width=Inches(2))

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()

def _make_xlsx_bytes() -> bytes:
    wb = Workbook()
    ws = cast(Worksheet, wb.active)
    ws.title = 'Alpha'
    ws.append(['Name', 'Score'])
    ws.append(['Alice', 90])
    ws2 = wb.create_sheet('Beta')
    ws2.append(['Item', 'Qty'])
    ws2.append(['Pen', 3])

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()

def _make_csv_bytes() -> bytes:
    return b'name,score\nAlice,95\nBob,88\n'

def _make_tsv_bytes() -> bytes:
    return 'name\tscore\nAlice\t95\nBob\t88\n'.encode('utf-8')

class TestFileField(TestCase):
    @classmethod
    def setUpClass(cls):
        cls.audio_wav_bytes = _make_audio_wav_bytes()
        cls.image_png_bytes = _make_image_png_bytes()
        cls.video_mp4_bytes = _make_video_mp4_bytes(cls.image_png_bytes, cls.audio_wav_bytes)
        cls.pdf_bytes = _make_pdf_bytes()
        cls.pdf_mixed_bytes = _make_pdf_mixed_bytes(cls.image_png_bytes)
        cls.docx_bytes = _make_docx_bytes(cls.image_png_bytes)
        cls.pptx_bytes = _make_pptx_bytes(cls.image_png_bytes)
        cls.xlsx_bytes = _make_xlsx_bytes()
        cls.html_bytes = (
            '<html><body><h1>HTML</h1><p>hello world</p><img src="data:image/png;base64,'
            + base64.b64encode(cls.image_png_bytes).decode('utf-8')
            + '"></body></html>'
        ).encode('utf-8')
        cls.rtf_bytes = r'{\rtf1\ansi RTF_TEXT}'.encode('utf-8')

        cls.audio = Audio(cls.audio_wav_bytes)
        cls.image = Image(cls.image_png_bytes)
        cls.video = Video(cls.video_mp4_bytes)
        cls.pdf = PDF.Load(cls.pdf_bytes)
        cls.pdf_mixed = PDF.Load(cls.pdf_mixed_bytes)
        cls.ppt = PPT.Load(cls.pptx_bytes)
        cls.html = HTML.Load(cls.html_bytes)
        cls.doc = Doc.Load(cls.docx_bytes)
        cls.rtf = RTF.Load(cls.rtf_bytes)
        cls.excel = Excel.Load(cls.xlsx_bytes)
        cls.json_bytes = b"{'name': 'demo', // comment\n 'items': [1,2]}"
        cls.yaml_bytes = b'name: demo\nitems:\n  - 1\n  - 2\n'
        cls.toml_bytes = b'title = "demo"\n[owner]\nname = "Bot"\n'
        cls.xml_bytes = b'<root><item>A</item><item>B</item></root>'
        cls.csv_bytes = _make_csv_bytes()
        cls.tsv_bytes = _make_tsv_bytes()
        cls.txt_bytes = b'TXT CONTENT\nNEXT LINE\n'
        cls.plaintext_bytes = b'PLAINTEXT\nNEXT LINE\n'
        cls.markdown_bytes = b'# Title\n\n- one\n- two\n'

        cls.json_doc = JSON.Load(cls.json_bytes)
        cls.yaml_doc = YAML.Load(cls.yaml_bytes)
        cls.toml_doc = TOML.Load(cls.toml_bytes)
        cls.xml_doc = XML.Load(cls.xml_bytes)
        cls.csv_doc = CSV.Load(cls.csv_bytes)
        cls.tsv_doc = TSV.Load(cls.tsv_bytes)
        cls.txt_doc = TXT.Load(cls.txt_bytes)
        cls.plaintext_doc = PlainText.Load(cls.plaintext_bytes)
        cls.markdown_doc = Markdown.Load(cls.markdown_bytes)
        cls.legacy_doc_path = _TEST_RESOURCES_DIR / 'legacy_sample_u8issue2.doc'
        cls.legacy_doc = Doc.Load(cls.legacy_doc_path) if cls.legacy_doc_path.exists() else None

    def test_media_payload_model_dump_and_validate_json(self):
        model = MediaPayloadModel(
            audio=self.audio,
            image=self.image,
            video=self.video,
            pdf=self.pdf,
        )

        dumped_json = model.model_dump_json()
        dumped = json.loads(dumped_json)
        self.assertEqual(dumped['audio']['type'], 'audio')
        self.assertEqual(dumped['image']['type'], 'image')
        self.assertEqual(dumped['video']['type'], 'video')
        self.assertEqual(dumped['pdf']['type'], 'pdf')

        roundtrip = MediaPayloadModel.model_validate_json(dumped_json)
        self.assertIsInstance(roundtrip.audio, Audio)
        self.assertIsInstance(roundtrip.image, Image)
        self.assertIsInstance(roundtrip.video, Video)
        self.assertIsInstance(roundtrip.pdf, PDF)
        
    def test_media_payload_model_validate_from_dict(self):
        payload = {
            'audio': {'type': 'audio', 'data': self.audio_wav_bytes},
            'image': {'type': 'image', 'data': self.image_png_bytes},
            'video': {'type': 'video', 'data': self.video_mp4_bytes},
            'pdf': {'type': 'pdf', 'data': self.pdf_bytes},
        }
        model = MediaPayloadModel.model_validate(payload)
        self.assertIsInstance(model.audio, Audio)
        self.assertIsInstance(model.image, Image)
        self.assertIsInstance(model.video, Video)
        self.assertIsInstance(model.pdf, PDF)

    def test_document_payload_model_dump_and_validate_json(self):
        model = DocumentPayloadModel(
            pdf=self.pdf,
            ppt=self.ppt,
            html=self.html,
            doc=self.doc,
            rtf=self.rtf,
            excel=self.excel,
        )

        dumped_json = model.model_dump_json()
        dumped = json.loads(dumped_json)
        self.assertEqual(dumped['pdf']['type'], 'pdf')
        self.assertEqual(dumped['ppt']['type'], 'ppt')
        self.assertEqual(dumped['html']['type'], 'html')
        self.assertEqual(dumped['doc']['type'], 'doc')
        self.assertEqual(dumped['rtf']['type'], 'rtf')
        self.assertEqual(dumped['excel']['type'], 'excel')

        roundtrip = DocumentPayloadModel.model_validate_json(dumped_json)
        self.assertIsInstance(roundtrip.pdf, PDF)
        self.assertIsInstance(roundtrip.ppt, PPT)
        self.assertIsInstance(roundtrip.html, HTML)
        self.assertIsInstance(roundtrip.doc, Doc)
        self.assertIsInstance(roundtrip.rtf, RTF)
        self.assertIsInstance(roundtrip.excel, Excel)

    def test_document_payload_model_validate_from_dict(self):
        payload = {
            'pdf': {'type': 'pdf', 'data': self.pdf_bytes},
            'ppt': {'type': 'ppt', 'data': self.pptx_bytes},
            'html': {'type': 'html', 'data': self.html_bytes},
            'doc': {'type': 'doc', 'data': self.docx_bytes},
            'rtf': {'type': 'rtf', 'data': self.rtf_bytes},
            'excel': {'type': 'excel', 'data': self.xlsx_bytes},
        }
        model = DocumentPayloadModel.model_validate(payload)
        self.assertIsInstance(model.pdf, PDF)
        self.assertIsInstance(model.ppt, PPT)
        self.assertIsInstance(model.html, HTML)
        self.assertIsInstance(model.doc, Doc)
        self.assertIsInstance(model.rtf, RTF)
        self.assertIsInstance(model.excel, Excel)

    def test_structured_document_payload_model_dump_and_validate_json(self):
        model = StructuredDocumentPayloadModel(
            json_doc=self.json_doc,
            yaml_doc=self.yaml_doc,
            toml_doc=self.toml_doc,
            xml_doc=self.xml_doc,
            csv_doc=self.csv_doc,
            tsv_doc=self.tsv_doc,
            txt_doc=self.txt_doc,
            plaintext_doc=self.plaintext_doc,
            markdown_doc=self.markdown_doc,
        )
        dumped = json.loads(model.model_dump_json())
        self.assertEqual(dumped['json_doc']['type'], 'json')
        self.assertEqual(dumped['yaml_doc']['type'], 'yaml')
        self.assertEqual(dumped['toml_doc']['type'], 'toml')
        self.assertEqual(dumped['xml_doc']['type'], 'xml')
        self.assertEqual(dumped['csv_doc']['type'], 'csv')
        self.assertEqual(dumped['tsv_doc']['type'], 'tsv')
        self.assertEqual(dumped['txt_doc']['type'], 'txt')
        self.assertEqual(dumped['plaintext_doc']['type'], 'plaintext')
        self.assertEqual(dumped['markdown_doc']['type'], 'markdown')

        roundtrip = StructuredDocumentPayloadModel.model_validate_json(model.model_dump_json())
        self.assertIsInstance(roundtrip.json_doc, JSON)
        self.assertIsInstance(roundtrip.yaml_doc, YAML)
        self.assertIsInstance(roundtrip.toml_doc, TOML)
        self.assertIsInstance(roundtrip.xml_doc, XML)
        self.assertIsInstance(roundtrip.csv_doc, CSV)
        self.assertIsInstance(roundtrip.tsv_doc, TSV)
        self.assertIsInstance(roundtrip.txt_doc, TXT)
        self.assertIsInstance(roundtrip.plaintext_doc, PlainText)
        self.assertIsInstance(roundtrip.markdown_doc, Markdown)

    def test_structured_document_payload_model_validate_from_dict(self):
        payload = {
            'json_doc': {'type': 'json', 'data': self.json_bytes},
            'yaml_doc': {'type': 'yaml', 'data': self.yaml_bytes},
            'toml_doc': {'type': 'toml', 'data': self.toml_bytes},
            'xml_doc': {'type': 'xml', 'data': self.xml_bytes},
            'csv_doc': {'type': 'csv', 'data': self.csv_bytes},
            'tsv_doc': {'type': 'tsv', 'data': self.tsv_bytes},
            'txt_doc': {'type': 'txt', 'data': self.txt_bytes},
            'plaintext_doc': {'type': 'plaintext', 'data': self.plaintext_bytes},
            'markdown_doc': {'type': 'markdown', 'data': self.markdown_bytes},
        }
        model = StructuredDocumentPayloadModel.model_validate(payload)
        self.assertIsInstance(model.json_doc, JSON)
        self.assertIsInstance(model.yaml_doc, YAML)
        self.assertIsInstance(model.toml_doc, TOML)
        self.assertIsInstance(model.xml_doc, XML)
        self.assertIsInstance(model.csv_doc, CSV)
        self.assertIsInstance(model.tsv_doc, TSV)
        self.assertIsInstance(model.txt_doc, TXT)
        self.assertIsInstance(model.plaintext_doc, PlainText)
        self.assertIsInstance(model.markdown_doc, Markdown)
        
    def test_file_payload_model_with_each_file_instance(self):
        for item, expected_type in (
            (self.audio, 'audio'),
            (self.image, 'image'),
            (self.video, 'video'),
            (self.pdf, 'pdf'),
            (self.ppt, 'ppt'),
            (self.html, 'html'),
            (self.doc, 'doc'),
            (self.rtf, 'rtf'),
            (self.excel, 'excel'),
            (self.json_doc, 'json'),
            (self.yaml_doc, 'yaml'),
            (self.toml_doc, 'toml'),
            (self.xml_doc, 'xml'),
            (self.csv_doc, 'csv'),
            (self.tsv_doc, 'tsv'),
            (self.txt_doc, 'txt'),
            (self.plaintext_doc, 'plaintext'),
            (self.markdown_doc, 'markdown'),
        ):
            model = FilePayloadModel(data=item)    # type: ignore[arg-type]
            dumped_json = model.model_dump_json()
            dumped = json.loads(dumped_json)
            self.assertEqual(dumped['data']['type'], expected_type)
            roundtrip = FilePayloadModel.model_validate_json(dumped_json)
            if expected_type == 'audio':
                self.assertIsInstance(roundtrip.data, Audio)
            elif expected_type == 'image':
                self.assertIsInstance(roundtrip.data, Image)
            elif expected_type == 'video':
                self.assertIsInstance(roundtrip.data, Video)
            elif expected_type == 'pdf':
                self.assertIsInstance(roundtrip.data, PDF)
            elif expected_type == 'ppt':
                self.assertIsInstance(roundtrip.data, PPT)
            elif expected_type == 'html':
                self.assertIsInstance(roundtrip.data, HTML)
            elif expected_type == 'doc':
                self.assertIsInstance(roundtrip.data, Doc)
            elif expected_type == 'rtf':
                self.assertIsInstance(roundtrip.data, RTF)
            elif expected_type == 'excel':
                self.assertIsInstance(roundtrip.data, Excel)
            elif expected_type == 'json':
                self.assertIsInstance(roundtrip.data, JSON)
            elif expected_type == 'yaml':
                self.assertIsInstance(roundtrip.data, YAML)
            elif expected_type == 'toml':
                self.assertIsInstance(roundtrip.data, TOML)
            elif expected_type == 'xml':
                self.assertIsInstance(roundtrip.data, XML)
            elif expected_type == 'csv':
                self.assertIsInstance(roundtrip.data, CSV)
            elif expected_type == 'tsv':
                self.assertIsInstance(roundtrip.data, TSV)
            elif expected_type == 'txt':
                self.assertIsInstance(roundtrip.data, TXT)
            elif expected_type == 'plaintext':
                self.assertIsInstance(roundtrip.data, PlainText)
            elif expected_type == 'markdown':
                self.assertIsInstance(roundtrip.data, Markdown)

    def test_legacy_doc_resource_roundtrip(self):
        self.assertIsNotNone(self.legacy_doc)
        legacy_doc = cast(Doc, self.legacy_doc)
        model = FilePayloadModel(data=legacy_doc)
        dumped_json = model.model_dump_json()
        roundtrip = FilePayloadModel.model_validate_json(dumped_json)
        self.assertIsInstance(roundtrip.data, Doc)
        self.assertEqual(roundtrip.data.to_bytes(), self.legacy_doc_path.read_bytes())
        
    def test_file_payload_model_validate_from_dict(self):
        cases = (
            ({'type': 'audio', 'data': self.audio_wav_bytes}, 'audio'),
            ({'type': 'image', 'data': self.image_png_bytes}, 'image'),
            ({'type': 'video', 'data': self.video_mp4_bytes}, 'video'),
            ({'type': 'pdf', 'data': self.pdf_bytes}, 'pdf'),
            ({'type': 'ppt', 'data': self.pptx_bytes}, 'ppt'),
            ({'type': 'html', 'data': self.html_bytes}, 'html'),
            ({'type': 'doc', 'data': self.docx_bytes}, 'doc'),
            ({'type': 'rtf', 'data': self.rtf_bytes}, 'rtf'),
            ({'type': 'excel', 'data': self.xlsx_bytes}, 'excel'),
            ({'type': 'json', 'data': self.json_bytes}, 'json'),
            ({'type': 'yaml', 'data': self.yaml_bytes}, 'yaml'),
            ({'type': 'toml', 'data': self.toml_bytes}, 'toml'),
            ({'type': 'xml', 'data': self.xml_bytes}, 'xml'),
            ({'type': 'csv', 'data': self.csv_bytes}, 'csv'),
            ({'type': 'tsv', 'data': self.tsv_bytes}, 'tsv'),
            ({'type': 'txt', 'data': self.txt_bytes}, 'txt'),
            ({'type': 'plaintext', 'data': self.plaintext_bytes}, 'plaintext'),
            ({'type': 'markdown', 'data': self.markdown_bytes}, 'markdown'),
        )
        for media_dict, expected_type in cases:
            model = FilePayloadModel.model_validate({'data': media_dict})
            if expected_type == 'audio':
                self.assertIsInstance(model.data, Audio)
            elif expected_type == 'image':
                self.assertIsInstance(model.data, Image)
            elif expected_type == 'video':
                self.assertIsInstance(model.data, Video)
            elif expected_type == 'pdf':
                self.assertIsInstance(model.data, PDF)
            elif expected_type == 'ppt':
                self.assertIsInstance(model.data, PPT)
            elif expected_type == 'html':
                self.assertIsInstance(model.data, HTML)
            elif expected_type == 'doc':
                self.assertIsInstance(model.data, Doc)
            elif expected_type == 'rtf':
                self.assertIsInstance(model.data, RTF)
            elif expected_type == 'excel':
                self.assertIsInstance(model.data, Excel)
            elif expected_type == 'json':
                self.assertIsInstance(model.data, JSON)
            elif expected_type == 'yaml':
                self.assertIsInstance(model.data, YAML)
            elif expected_type == 'toml':
                self.assertIsInstance(model.data, TOML)
            elif expected_type == 'xml':
                self.assertIsInstance(model.data, XML)
            elif expected_type == 'csv':
                self.assertIsInstance(model.data, CSV)
            elif expected_type == 'tsv':
                self.assertIsInstance(model.data, TSV)
            elif expected_type == 'txt':
                self.assertIsInstance(model.data, TXT)
            elif expected_type == 'plaintext':
                self.assertIsInstance(model.data, PlainText)
            elif expected_type == 'markdown':
                self.assertIsInstance(model.data, Markdown)

    def test_file_classes_expose_type_metadata(self):
        self.assertIsNotNone(cast(type[FileTypeMetaProtocol], Audio).Type)
        self.assertEqual('audio', Audio.Type)
        self.assertIn('mp3', Audio.TypeNames)
        self.assertEqual('image', Image.Type)
        self.assertIn('png', Image.TypeNames)
        self.assertEqual('video', Video.Type)
        self.assertIn('mp4', Video.TypeNames)
        self.assertEqual('pdf', PDF.Type)
        self.assertEqual('doc', Doc.Type)
        self.assertEqual('excel', Excel.Type)
        self.assertEqual('json', JSON.Type)
        self.assertEqual('yaml', YAML.Type)
        self.assertEqual('toml', TOML.Type)
        self.assertEqual('xml', XML.Type)
        self.assertEqual('csv', CSV.Type)
        self.assertEqual('tsv', TSV.Type)
        self.assertEqual('txt', TXT.Type)
        self.assertEqual('plaintext', PlainText.Type)
        self.assertEqual('markdown', Markdown.Type)

    def test_pdf_extract_page_contents_with_mixed_content(self):
        pages = self.pdf_mixed.extract_page_contents()
        self.assertIsInstance(pages, list)
        self.assertGreaterEqual(len(pages), 1)

        first_page = pages[0]
        self.assertIsInstance(first_page, list)
        self.assertGreaterEqual(len(first_page), 3)

        self.assertIsInstance(first_page[0], str)
        self.assertIn('MIX_TEXT_1', first_page[0])   # type: ignore
        self.assertIsInstance(first_page[1], Image)
        self.assertIsInstance(first_page[2], str)
        self.assertIn('MIX_TEXT_2', first_page[2])   # type: ignore
        
if __name__ == '__main__':
    main()